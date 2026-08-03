"""Pure PPTX-generation service: aggregates sales/stock/GST data from the DB and
builds a slide deck. No Telegram or agent imports here — see the standalone test
script for exercising this directly without the chat loop.
"""

import os
import tempfile
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone

import matplotlib

matplotlib.use("Agg")  # headless: no display available, only ever rendering to file
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import func, select

from app.db.models import Bill, BillItem, Product
from app.db.session import async_session_maker


@dataclass
class AnalysisData:
    date_from: str
    date_to: str
    bill_count: int
    total_sales: float
    gst_collected: float
    top_products: list[tuple[str, float]]
    stock_levels: list[tuple[str, float]]


def _parse_range(date_from: str, date_to: str) -> tuple[datetime, datetime]:
    try:
        start = datetime.strptime(date_from, "%Y-%m-%d").replace(tzinfo=timezone.utc)
        end = datetime.strptime(date_to, "%Y-%m-%d").replace(tzinfo=timezone.utc)
    except ValueError as exc:
        raise ValueError(f"date_from/date_to must be YYYY-MM-DD: {exc}") from exc
    end = end + timedelta(days=1, microseconds=-1)  # inclusive end of day
    if start > end:
        raise ValueError(f"date_from ({date_from}) is after date_to ({date_to}).")
    return start, end


async def _aggregate(shop_id: str, date_from: str, date_to: str) -> AnalysisData:
    start, end = _parse_range(date_from, date_to)

    async with async_session_maker() as session:
        bills_stmt = select(Bill).where(
            Bill.shop_id == shop_id,
            Bill.status == "finalized",
            Bill.finalized_at >= start,
            Bill.finalized_at <= end,
        )
        bills = list((await session.execute(bills_stmt)).scalars().all())
        bill_ids = [b.id for b in bills]

        total_sales = sum(float(b.grand_total) for b in bills)
        gst_collected = sum(float(b.cgst_total) + float(b.sgst_total) for b in bills)

        top_products: list[tuple[str, float]] = []
        if bill_ids:
            items_stmt = (
                select(Product.name, func.sum(BillItem.quantity))
                .join(Product, Product.id == BillItem.product_id)
                .where(BillItem.bill_id.in_(bill_ids))
                .group_by(Product.name)
                .order_by(func.sum(BillItem.quantity).desc())
                .limit(8)
            )
            rows = (await session.execute(items_stmt)).all()
            top_products = [(name, float(qty)) for name, qty in rows]

        stock_stmt = (
            select(Product.name, Product.quantity_on_hand)
            .where(Product.shop_id == shop_id)
            .order_by(Product.quantity_on_hand.desc())
            .limit(10)
        )
        stock_rows = (await session.execute(stock_stmt)).all()
        stock_levels = [(name, float(qty)) for name, qty in stock_rows]

    return AnalysisData(
        date_from=date_from,
        date_to=date_to,
        bill_count=len(bills),
        total_sales=total_sales,
        gst_collected=gst_collected,
        top_products=top_products,
        stock_levels=stock_levels,
    )


def _bar_chart_png(labels: list[str], values: list[float], title: str, ylabel: str) -> str:
    fig, ax = plt.subplots(figsize=(9, 4.8))
    if labels:
        ax.bar(labels, values, color="#4C72B0")
    else:
        ax.text(0.5, 0.5, "No data in range", ha="center", va="center", transform=ax.transAxes)
    ax.set_title(title)
    ax.set_ylabel(ylabel)
    ax.tick_params(axis="x", rotation=30)
    fig.tight_layout()

    fd, path = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    fig.savefig(path, dpi=150)
    plt.close(fig)
    return path


def _add_summary_slide(prs: Presentation, data: AnalysisData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = (
        f"Sales & Inventory Analysis: {data.date_from} to {data.date_to}"
    )

    rows = [
        ("Finalized bills", str(data.bill_count)),
        ("Total sales (grand total)", f"{data.total_sales:.2f}"),
        ("GST collected (CGST+SGST)", f"{data.gst_collected:.2f}"),
    ]
    table_shape = slide.shapes.add_table(
        rows=len(rows) + 1, cols=2, left=Inches(0.7), top=Inches(1.6), width=Inches(8.5), height=Inches(2.2)
    )
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    for i, (label, value) in enumerate(rows, start=1):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = value
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(16)


def _add_chart_slide(prs: Presentation, title: str, png_path: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(png_path, Inches(0.6), Inches(1.5), width=Inches(8.8))


async def generate_analysis_pptx(shop_id: str, date_from: str, date_to: str) -> str:
    """Build a sales/stock/GST analysis deck for [date_from, date_to] (inclusive,
    both "YYYY-MM-DD") and return the path to the generated .pptx file. Raises
    ValueError for a malformed or inverted date range."""
    data = await _aggregate(shop_id, date_from, date_to)

    prs = Presentation()
    _add_summary_slide(prs, data)

    top_labels = [name for name, _ in data.top_products]
    top_values = [qty for _, qty in data.top_products]
    top_png = _bar_chart_png(top_labels, top_values, "Top Products by Quantity Sold", "Quantity")
    try:
        _add_chart_slide(prs, "Top Products by Quantity Sold", top_png)
    finally:
        os.remove(top_png)

    stock_labels = [name for name, _ in data.stock_levels]
    stock_values = [qty for _, qty in data.stock_levels]
    stock_png = _bar_chart_png(stock_labels, stock_values, "Current Stock Levels", "Quantity on Hand")
    try:
        _add_chart_slide(prs, "Current Stock Levels", stock_png)
    finally:
        os.remove(stock_png)

    fd, path = tempfile.mkstemp(prefix="sales_analysis_", suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path
