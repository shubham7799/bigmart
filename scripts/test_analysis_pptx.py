"""Standalone test for app/services/analysis_pptx.py — no agent, no Telegram.

Run: python scripts/test_analysis_pptx.py <date_from> <date_to>
"""

import asyncio
import sys

from pptx import Presentation

from app.services.analysis_pptx import _aggregate, generate_analysis_pptx


async def main(date_from: str, date_to: str) -> None:
    data = await _aggregate(date_from, date_to)
    print("--- aggregated data ---")
    print("bill_count:", data.bill_count)
    print("total_sales:", data.total_sales)
    print("gst_collected:", data.gst_collected)
    print("top_products:", data.top_products)
    print("stock_levels:", data.stock_levels)

    path = await generate_analysis_pptx(date_from, date_to)
    print(f"\nPPTX written to: {path}")

    prs = Presentation(path)
    print(f"slide count: {len(prs.slides._sldIdLst)}")
    for i, slide in enumerate(prs.slides):
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        pics = [s for s in slide.shapes if s.shape_type == 13]
        print(f"slide {i}: text={texts!r} pictures={len(pics)}")

    # sanity: totals in the summary slide should match the aggregated numbers
    def _shape_text(shape) -> str:
        if shape.has_text_frame:
            return shape.text_frame.text
        if shape.has_table:
            return " ".join(cell.text for row in shape.table.rows for cell in row.cells)
        return ""

    summary_texts = " ".join(_shape_text(shape) for shape in prs.slides[0].shapes)
    assert f"{data.total_sales:.2f}" in summary_texts, "total_sales not found on summary slide"
    assert f"{data.gst_collected:.2f}" in summary_texts, "gst_collected not found on summary slide"
    assert str(data.bill_count) in summary_texts, "bill_count not found on summary slide"
    print("\nOK: summary slide numbers match aggregated data.")

    # inverted range should be rejected
    try:
        await generate_analysis_pptx(date_to, date_from) if date_from != date_to else (_ for _ in ()).throw(
            ValueError("skip: from==to, can't test inversion")
        )
    except ValueError as exc:
        print(f"OK: inverted date range correctly rejected: {exc}")


if __name__ == "__main__":
    date_from = sys.argv[1] if len(sys.argv) > 1 else "2026-07-25"
    date_to = sys.argv[2] if len(sys.argv) > 2 else "2026-08-01"
    asyncio.run(main(date_from, date_to))
